import streamlit as st
from pptx import Presentation
from langchain_community.document_loaders import DataFrameLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.llms import HuggingFaceHub
from langchain.chains import RetrievalQA
import pandas as pd
import io

# Configuración de la página
st.set_page_config(page_title="Legal Check Plus AI", page_icon="⚖️")
st.title("⚖️ Asistente Legal Check Plus")

# --- BARRA LATERAL PARA CARGA DE ARCHIVOS ---
with st.sidebar:
    st.header("Configuración")
    uploaded_file = st.file_uploader("Sube la presentación comercial (.pptx)", type="pptx")
    st.info("El archivo se procesa en memoria RAM y no se guarda en ningún servidor.")

# --- FUNCIÓN PARA EXTRAER TEXTO DEL PPTX ---
def extract_text_from_pptx(file):
    prs = Presentation(file)
    text_data = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        full_text = " ".join(slide_text)
        if full_text:
            text_data.append({"text": full_text, "source": f"Slide {i+1}"})
    return pd.DataFrame(text_data)

# --- PROCESAMIENTO DEL AGENTE ---
if uploaded_file is not None:
    @st.cache_resource
    def initialize_agent(_file):
        # 1. Extracción desde el archivo subido
        df = extract_text_from_pptx(_file)
        loader = DataFrameLoader(df, page_content_column="text")
        documents = loader.load()

        # 2. División de texto
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        texts = text_splitter.split_documents(documents)

        # 3. Embeddings Gratuitos
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        
        # 4. Vector Store
        vectorstore = FAISS.from_documents(texts, embeddings)
        
        # 5. LLM (Mistral-7B)
        llm = HuggingFaceHub(
            repo_id="mistralai/Mistral-7B-Instruct-v0.2",
            model_kwargs={"temperature": 0.1, "max_length": 512}
        )

        # 6. Cadena de Respuesta
        return RetrievalQA.from_chain_type(
            llm=llm,
            chain_type="stuff",
            retriever=vectorstore.as_retriever()
        )

    try:
        agent = initialize_agent(uploaded_file)
        st.success("✅ Documento cargado exitosamente. ¡Ya puedes preguntar!")
    except Exception as e:
        st.error(f"Error procesando el archivo: {e}")
        st.stop()

    # --- INTERFAZ DE CHAT ---
    if "messages" not in st.session_state:
        st.session_state.messages = []

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    if prompt := st.chat_input("¿Qué deseas saber sobre Legal Check Plus?"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        with st.chat_message("assistant"):
            with st.spinner("Analizando presentación..."):
                # Prompt optimizado para español y fidelidad al documento
                query = f"Responde en español basándote estrictamente en el documento cargado. Si la información no está, di que no lo sabes. Pregunta: {prompt}"
                response = agent.run(query)
                st.markdown(response)
                st.session_state.messages.append({"role": "assistant", "content": response})
else:
    st.warning("⚠️ Por favor, sube el archivo .pptx en la barra lateral para comenzar.")
